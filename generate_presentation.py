"""
Generate a PowerPoint presentation for the Roofing Lead Intelligence Platform.
Interview-focused: structured for a 1-hour Q&A session with interviewers.
Includes the agentic chat feature.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BLACK = RGBColor(0x17, 0x17, 0x17)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
MED_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFE, 0xFF, 0xF7)
ORANGE = RGBColor(0xEA, 0x88, 0x0E)
TEAL = RGBColor(0x0D, 0x9B, 0x8C)
GREEN = RGBColor(0x16, 0x83, 0x5E)
RED = RGBColor(0xAA, 0x1E, 0x00)
BLUE = RGBColor(0x1E, 0x40, 0xAF)

FONT_BODY = "Calibri"
FONT_MONO = "Consolas"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, font_name=FONT_BODY,
                alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    return txBox


def add_subtitle(slide, text="Roofing Lead Intelligence Platform \u2014 Aryaman Gupta | Instalily", top=Inches(0.35)):
    add_textbox(slide, Inches(0.8), top, Inches(11.7), Inches(0.35),
                text, font_size=11, color=MED_GRAY)


def add_title(slide, text, top=Inches(0.9)):
    add_textbox(slide, Inches(0.8), top, Inches(11.7), Inches(0.6),
                text, font_size=36, bold=True, color=BLACK)


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, left, top, width, height, items, font_size=16,
                color=BLACK, bullet="\u2022  ", spacing=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}{item}" if bullet else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)
    return txBox


# ══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
BL = prs.slide_layouts[6]  # Blank


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, BLACK)
add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
            "Roofing Lead Intelligence Platform",
            font_size=44, bold=True, color=WHITE)
add_textbox(s, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.8),
            "AI-powered sales intelligence for roofing distributors",
            font_size=24, color=LIGHT_GRAY)
add_textbox(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.6),
            "Scrape \u2192 Research \u2192 Score \u2192 Extract Contacts \u2192 Dashboard + Agentic Chat",
            font_size=16, color=MED_GRAY, font_name=FONT_MONO)
add_textbox(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
            "Aryaman Gupta  \u00b7  Onsite Project  \u00b7  Instalily AI",
            font_size=16, color=MED_GRAY)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "The Problem")
add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), [
    "Roofing distributors have no systematic way to identify and prioritize contractor leads",
    "Sales reps spend hours manually researching contractors \u2014 Googling names, checking reviews, guessing who to call first",
    "GAF\u2019s public directory has rich contractor data, but it\u2019s locked inside a JS-rendered, bot-protected website",
    "No existing tool combines scraping + AI enrichment + actionable sales intelligence in one workflow",
    "Without lead scoring, reps waste time on low-potential contractors while high-value leads go untouched",
], font_size=20, spacing=14)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE SOLUTION (OVERVIEW)
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "The Solution")
add_textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.7),
            "An end-to-end platform that scrapes, enriches, scores, and presents contractor leads \u2014 plus an agentic chat interface for natural-language exploration.",
            font_size=18, color=DARK_GRAY)

cards = [
    ("Scrape", "Playwright + Coveo API\ninterception", "78 contractors from\nGAF directory"),
    ("Research", "Perplexity sonar-pro\nweb research engine", "Grounded insights with\ncitations per lead"),
    ("Score + Enrich", "Hybrid deterministic +\nLLM scoring (0\u2013100)", "5-factor breakdown +\ntalking points, emails"),
    ("Dashboard", "Next.js + FastAPI\nInstalily-branded UI", "Filter, sort, drill into\nevery lead"),
    ("Agent Chat", "OpenAI function calling\nagentic loop (4 tools)", "Natural language queries\nover the leads DB"),
]
cw, ch, gap = Inches(2.2), Inches(3.0), Inches(0.2)
sx = Inches(0.5)
cy = Inches(2.8)
for i, (title, desc, detail) in enumerate(cards):
    x = sx + i * (cw + gap)
    add_rect(s, x, cy, cw, ch, WHITE, BLACK)
    add_textbox(s, x + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), Inches(0.5),
                title, font_size=20, bold=True)
    add_textbox(s, x + Inches(0.2), cy + Inches(0.85), cw - Inches(0.4), Inches(0.9),
                desc, font_size=14, color=DARK_GRAY)
    add_textbox(s, x + Inches(0.2), cy + Inches(1.85), cw - Inches(0.4), Inches(0.9),
                detail, font_size=13, color=MED_GRAY, font_name=FONT_MONO)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 4 — ARCHITECTURE DIAGRAM
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Architecture: Pipeline + API + Dashboard + Agent")

flow = [
    ("GAF\nWebsite", DARK_GRAY),
    ("Playwright\nScraper", TEAL),
    ("Perplexity\nResearch", BLUE),
    ("OpenAI\nScoring +\nContacts", ORANGE),
    ("SQLite\nDatabase", GREEN),
    ("FastAPI\nREST API", DARK_GRAY),
    ("Next.js\nDashboard", BLACK),
]
bw, bh, aw = Inches(1.45), Inches(1.15), Inches(0.3)
bx, by = Inches(0.4), Inches(2.2)
for i, (label, color) in enumerate(flow):
    x = bx + i * (bw + aw)
    sh = add_rect(s, x, by, bw, bh, color)
    sh.text_frame.paragraphs[0].text = label
    sh.text_frame.paragraphs[0].font.size = Pt(12)
    sh.text_frame.paragraphs[0].font.color.rgb = WHITE
    sh.text_frame.paragraphs[0].font.bold = True
    sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sh.text_frame.word_wrap = True
    if i < len(flow) - 1:
        add_textbox(s, x + bw, by + Inches(0.25), aw, Inches(0.6),
                    "\u2192", font_size=26, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Agent chat branch
add_textbox(s, Inches(7.7), Inches(3.55), Inches(1.2), Inches(0.4),
            "\u2191", font_size=26, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
sh = add_rect(s, Inches(6.7), Inches(3.9), Inches(2.8), Inches(0.8), ORANGE)
sh.text_frame.paragraphs[0].text = "Agentic Chat (OpenAI function calling)"
sh.text_frame.paragraphs[0].font.size = Pt(12)
sh.text_frame.paragraphs[0].font.color.rgb = WHITE
sh.text_frame.paragraphs[0].font.bold = True
sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
sh.text_frame.word_wrap = True

add_bullets(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.5), [
    "Pipeline runs offline \u2014 dashboard and chat never wait on scraping or LLM enrichment calls",
    "Each pipeline stage is independently runnable: re-score without re-researching, re-research without re-scraping",
    "Chat agent queries the same DB via FastAPI \u2014 same service layer, same filters, same data",
    "Clean separation: CLI for pipeline ops, REST API for reads, Next.js for presentation, Chat for exploration",
], font_size=15, spacing=9)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 5 — SCRAPING
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Scraping: Playwright + Coveo API Interception")
add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2), [
    "GAF\u2019s site is JS-rendered and protected by Akamai bot detection \u2014 simple HTTP requests won\u2019t work",
    "Playwright launches headless Chromium with stealth settings (spoofed webdriver, plugins, permissions)",
    "Key insight: Instead of fragile DOM parsing, we intercept the Coveo search-as-a-service API calls for clean JSON",
    "Extract Bearer token from the first request, then paginate directly against Coveo\u2019s REST API",
    "Profile page scraping (configurable concurrency) extracts website URLs, years in business, about text",
    "Parameterized: zip_code and distance are inputs \u2014 works for any territory, not hardcoded to 10013",
    "Result: 78 contractors with name, address, phone, certification, rating, reviews, services, coordinates, distance",
], font_size=18, spacing=12)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 6 — THREE-STAGE AI ENRICHMENT
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Three-Stage AI Enrichment Pipeline")
add_textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5),
            "Multiple LLMs, each doing what they\u2019re best at. Stages run independently with --force, --limit, --concurrency flags.",
            font_size=16, color=DARK_GRAY)

stages = [
    ("Stage 1: Perplexity \u2014 Web Research", BLUE, [
        "Queries Perplexity sonar-pro for each contractor: company overview, decision-makers, review sentiment, storm activity, growth signals, BBB, supplier intel",
        "Returns grounded citations (URLs) alongside research text \u2014 reps can verify every claim",
    ]),
    ("Stage 2: OpenAI \u2014 Hybrid Scoring + Insights", ORANGE, [
        "Deterministic (60/100 pts): certification tier (0\u201330), review volume log-scaled (0\u201320), rating quality (0\u201310)",
        "LLM-scored (40/100 pts): business signals (0\u201320) + why-now urgency (0\u201320) extracted from research text",
        "Also generates: talking points, buying signals, pain points, recommended pitch, draft outreach email",
    ]),
    ("Stage 3: OpenAI \u2014 Contact Extraction", GREEN, [
        "Extracts decision-maker contacts from research text: names, titles, emails, LinkedIn URLs",
        "Provider-based architecture: PerplexityExtractor (implemented), Hunter.io (stub, ready to plug in)",
    ]),
]
st = Inches(2.3)
for i, (title, color, bullets) in enumerate(stages):
    bh = Inches(0.33) * len(bullets) + Inches(0.55)
    y = st
    add_rect(s, Inches(0.8), y, Inches(0.08), bh, color)
    add_textbox(s, Inches(1.1), y + Inches(0.05), Inches(11.0), Inches(0.35),
                title, font_size=18, bold=True)
    for j, b in enumerate(bullets):
        add_textbox(s, Inches(1.3), y + Inches(0.42) + j * Inches(0.35),
                    Inches(10.8), Inches(0.35),
                    f"\u2022  {b}", font_size=14, color=DARK_GRAY)
    st += bh + Inches(0.1)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 7 — HYBRID SCORING DEEP DIVE
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Hybrid Scoring: Why Not Fully LLM or Fully Deterministic?")

cw = Inches(5.5)
# Left col
add_textbox(s, Inches(0.8), Inches(1.9), cw, Inches(0.4),
            "Deterministic (60 pts \u2014 Python)", font_size=20, bold=True, color=TEAL)
add_bullets(s, Inches(0.8), Inches(2.5), cw, Inches(2.5), [
    "certification_tier (0\u201330): President\u2019s Club=30, Master Elite=25, Certified=15",
    "review_volume (0\u201320): log-scaled from review count",
    "rating_quality (0\u201310): star rating weighted by review confidence",
    "Same input = same score. Costs $0. Fully auditable.",
], font_size=15, spacing=10)
# Right col
add_textbox(s, Inches(7.0), Inches(1.9), cw, Inches(0.4),
            "LLM-Scored (40 pts \u2014 gpt-4o-mini)", font_size=20, bold=True, color=ORANGE)
add_bullets(s, Inches(7.0), Inches(2.5), cw, Inches(2.5), [
    "business_signals (0\u201320): growth, hiring, expansion from research",
    "why_now_urgency (0\u201320): storms, seasonal demand, supplier issues",
    "Only an LLM can extract these from 3000 chars of unstructured text",
    "~$0.001 per contractor. All 78 leads scored for ~$0.05 total.",
], font_size=15, spacing=10)
# Benefits
add_textbox(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.4),
            "Why this matters:", font_size=18, bold=True)
add_bullets(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.8), [
    "Graceful degradation: if OpenAI is down, deterministic subtotal (60/100) still provides usable ranking",
    "Pre-computed deterministic scores passed to LLM as \u201clocked \u2014 do not modify\u201d \u2014 prevents hallucinated scores",
    "LLM scores clamped to valid ranges; final arithmetic enforced in Python (don\u2019t trust the LLM\u2019s addition)",
    "response_format={\"type\": \"json_object\"} guarantees parseable output every time",
], font_size=15, spacing=8)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 8 — AGENTIC CHAT
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Agentic Chat: Natural Language Over the Leads DB")

add_textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5),
            "Sales reps can ask questions in plain English. The LLM decides which tools to call, executes them, and iterates until it has an answer.",
            font_size=17, color=DARK_GRAY)

# Architecture diagram
add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.4),
            "User  \u2192  Next.js ChatPanel  \u2192  POST /api/chat  \u2192  Agent Loop (up to 5 rounds):  LLM decides \u2192 Tool call \u2192 Result \u2192 LLM decides \u2192 \u2026 \u2192 Final text",
            font_size=14, color=DARK_GRAY, font_name=FONT_MONO)

# 4 tools
tools = [
    ("search_leads", "Search & filter the leads database\nwith score, certification, name, sort, pagination"),
    ("get_lead_detail", "Full detail for one contractor:\ninsights, score breakdown, talking points, contacts"),
    ("get_stats", "Aggregate dashboard statistics:\ntotals, averages, cert breakdown, score distribution"),
    ("compare_leads", "Side-by-side comparison of 2\u20133 leads\non score, certification, signals, insights"),
]
tw, th, tg = Inches(2.7), Inches(1.8), Inches(0.3)
tx, ty = Inches(0.6), Inches(3.1)
for i, (name, desc) in enumerate(tools):
    x = tx + i * (tw + tg)
    add_rect(s, x, ty, tw, th, WHITE, BLACK)
    add_textbox(s, x + Inches(0.15), ty + Inches(0.15), tw - Inches(0.3), Inches(0.4),
                name, font_size=15, bold=True, font_name=FONT_MONO)
    add_textbox(s, x + Inches(0.15), ty + Inches(0.6), tw - Inches(0.3), Inches(1.0),
                desc, font_size=13, color=DARK_GRAY)

# Implementation details
add_bullets(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(2.0), [
    "OpenAI function calling with tool_choice=\"auto\" \u2014 LLM chains multiple tools in a single turn",
    "Domain-aware system prompt: knows GAF tiers, score components, product categories, lead temperature thresholds",
    "Tools dispatch to the same service layer as the REST API \u2014 one source of truth, no data divergence",
    "Floating chat bubble UI with markdown rendering, loading animation, conversation history",
    'Example: "Who are the top 5 Master Elite contractors near me?" \u2192 search_leads(certification="Master Elite", sort_by="lead_score", per_page=5)',
], font_size=15, spacing=9)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 9 — DATABASE DESIGN
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Data Management: Three-Table Schema")

tables = [
    ("contractors", "23 columns", [
        "Raw scraped data from GAF",
        "gaf_id as dedupe key (upserts)",
        "Name, address, phone, website",
        "Certification, rating, reviews",
        "Lat/lng, distance, services",
        "Years in business, about text",
    ]),
    ("lead_insights", "15 columns", [
        "AI-generated enrichments",
        "1:1 with contractor (UNIQUE FK)",
        "lead_score (0\u2013100)",
        "score_breakdown (5-factor JSON)",
        "Talking points, buying signals",
        "Why now, draft email, pitch",
    ]),
    ("contacts", "10 columns", [
        "Decision-maker contacts",
        "Many-to-one with contractor",
        "Name, title, email, phone",
        "LinkedIn URL",
        "Source + confidence level",
        "Deduped by name + contractor",
    ]),
]
tw, th, tg = Inches(3.6), Inches(3.5), Inches(0.4)
tx, ty = Inches(0.8), Inches(1.9)
for i, (name, cols, fields) in enumerate(tables):
    x = tx + i * (tw + tg)
    add_rect(s, x, ty, tw, th, WHITE, BLACK)
    add_textbox(s, x + Inches(0.2), ty + Inches(0.15), tw - Inches(0.4), Inches(0.35),
                name, font_size=18, bold=True, font_name=FONT_MONO)
    add_textbox(s, x + Inches(0.2), ty + Inches(0.5), tw - Inches(0.4), Inches(0.25),
                cols, font_size=12, color=MED_GRAY, font_name=FONT_MONO)
    for j, f in enumerate(fields):
        add_textbox(s, x + Inches(0.2), ty + Inches(0.85) + j * Inches(0.38),
                    tw - Inches(0.4), Inches(0.35),
                    f"\u2022  {f}", font_size=13, color=DARK_GRAY)
add_bullets(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.5), [
    "SQLAlchemy ORM \u2014 swap SQLite for Postgres by changing one connection string",
    "JSON fields stored as TEXT (SQLite-compatible), become native JSONB with GIN indexes in Postgres",
    "Three tables enable independent re-processing: re-score without re-researching, re-research without re-scraping",
], font_size=15, spacing=8)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 10 — REST API
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "REST API: 7 Endpoints, 11+ Schemas, 3-Tier Architecture")

endpoints = [
    ("GET", "/api/leads", "Paginated list with filters (score range, cert, search), multi-column sorting"),
    ("GET", "/api/leads/{id}", "Full detail: contractor + insights + contacts (eager-loaded joins)"),
    ("GET", "/api/stats", "Dashboard metrics: totals, averages, cert breakdown, score distribution"),
    ("POST", "/api/chat", "Agentic chat \u2014 natural language queries with OpenAI function calling"),
    ("POST", "/api/pipeline/scrape", "Trigger scraping for any zip_code + distance"),
    ("POST", "/api/pipeline/enrich", "Run all 3 enrichment stages sequentially"),
    ("GET", "/api/pipeline/status", "Pipeline health: how many awaiting research/scoring/contacts"),
]
et = Inches(1.9)
for i, (method, path, desc) in enumerate(endpoints):
    y = et + i * Inches(0.52)
    mc = GREEN if method == "GET" else ORANGE
    add_textbox(s, Inches(0.8), y, Inches(0.7), Inches(0.4),
                method, font_size=14, bold=True, color=mc, font_name=FONT_MONO)
    add_textbox(s, Inches(1.5), y, Inches(2.5), Inches(0.4),
                path, font_size=14, color=BLACK, font_name=FONT_MONO)
    add_textbox(s, Inches(4.2), y, Inches(8.5), Inches(0.4),
                desc, font_size=14, color=DARK_GRAY)

add_bullets(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.8), [
    "Schemas: 11 Pydantic models enforce type safety. Separate list vs. detail schemas avoid serializing multi-KB fields per row.",
    "Service layer: Isolates SQLAlchemy query logic. NULL-safe ordering, CASE WHEN bucketing for stats.",
    "All handlers sync def \u2014 enricher uses asyncio.run() internally; FastAPI runs sync handlers in a thread pool.",
    "Column whitelist prevents SQL injection via dynamic sort. per_page capped at 100.",
], font_size=14, spacing=7)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 11 — DASHBOARD UI
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Dashboard: Instalily-Branded UI")
add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), [
    "Instalily design language: IBM Plex Sans/Mono, off-white #FEFFF7, sharp black borders, monospace uppercase section eyebrows",
    "StatsBar: 4 metric cells (Total Leads, Top Leads, Avg Score, Enriched) in sharp-bordered grid",
    "FilterBar: Inline score-tier pills (Hot 60+ / Warm / Cold) + certification pills + debounced text search (300ms)",
    "LeadTable: Sortable columns (score, name, volume, distance) with paginated rows, ScorePill + CertBadge inline",
    "All filter/sort/page state synced to URL query params \u2014 every view is bookmarkable and shareable",
    "Lead detail page: 8 insight sections \u2014 header, contacts, why-now banner, talking points, score breakdown, buying signals / pain points, draft email, research summary",
    "Floating chat bubble (bottom-right) opens the agentic chat panel for natural-language queries",
    "26 components, zero external UI libraries \u2014 pure Tailwind CSS v4 + Next.js 16 App Router",
], font_size=17, spacing=11)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 12 — LEAD DETAIL PAGE (8 SECTIONS)
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Lead Detail Page: 8 Insight Sections")

sections = [
    ("LeadHeader", "Name, score pill, cert badge,\nmeta (city, phone, website, yrs)"),
    ("ContactCard", "Decision-maker cards with\nemail / call / LinkedIn buttons"),
    ("WhyNowBanner", "Time-sensitive urgency signal\n(orange left-border accent)"),
    ("TalkingPoints", "Numbered list (01, 02, 03)\nmonospace formatting"),
    ("ScoreBreakdown", "5 labeled progress bars\nshowing each scoring factor"),
    ("InsightsGrid", "Buying signals (green) +\npain points (red) two-column"),
    ("DraftEmail", "Dark #171717 card with\none-click copy button"),
    ("ResearchSummary", "Collapsible text + citation\nsource pills (clickable links)"),
]
cw, ch = Inches(2.7), Inches(1.8)
gx, gy = Inches(0.3), Inches(0.25)
sx, sy = Inches(0.6), Inches(1.9)
for i, (name, desc) in enumerate(sections):
    c, r = i % 4, i // 4
    x = sx + c * (cw + gx)
    y = sy + r * (ch + gy)
    add_rect(s, x, y, cw, ch, WHITE, BLACK)
    add_textbox(s, x + Inches(0.15), y + Inches(0.15), cw - Inches(0.3), Inches(0.35),
                name, font_size=16, bold=True, font_name=FONT_MONO)
    add_textbox(s, x + Inches(0.15), y + Inches(0.55), cw - Inches(0.3), Inches(1.0),
                desc, font_size=13, color=DARK_GRAY)

add_textbox(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
            "Lead detail is a server component \u2014 single fetch, data passed as props. Only CopyButton and ResearchSummary are client components.",
            font_size=14, color=MED_GRAY)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 13 — WHAT REPS ACTUALLY GET (CONCRETE EXAMPLE)
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "What Sales Reps Actually Get")
add_textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.4),
            "Example: Grapevine Pro \u2014 Score 67 (Hot), President\u2019s Club, 449 reviews, Iselin NJ",
            font_size=17, bold=True, color=DARK_GRAY)

add_bullets(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.8), [
    'Score Breakdown: Certification 30/30 | Reviews 18/20 | Rating 9/10 | Biz Signals 10/20 | Why Now 0/20',
    'Talking Point: "Mention their 449 Google reviews \u2014 they clearly value reputation"',
    'Talking Point: "Ask about their storm restoration pipeline and material lead times"',
    'Buying Signal: "High volume of storm damage repair work"',
    'Pain Point: "Current supplier has inconsistent delivery times (per Google reviews)"',
    'Why Now: "Major hailstorm hit their primary service area 3 weeks ago"',
    'Recommended Pitch: "Position as reliable, high-volume supplier for storm season..."',
    'Draft Email: Personalized outreach with subject line, ready to copy and send',
    'Key Contact: Frank N. \u2014 Owner \u2014 email, phone, LinkedIn (with confidence level)',
    'Chat: "Compare Grapevine Pro with Great American Roofing" \u2192 side-by-side analysis via agent',
], font_size=16, spacing=9)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 14 — TESTING
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Testing: 119 Tests Across the Full Backend")

cats = [
    ("Service Layer \u2014 25 tests", [
        "Pagination, all filter combinations, sort ordering, NULL score handling",
        "Stats aggregation, score distribution bucketing",
    ]),
    ("API Endpoints \u2014 19 tests", [
        "Full HTTP round-trip via FastAPI TestClient with in-memory SQLite (StaticPool)",
        "Response shapes, pagination metadata, 404 handling, input validation (422)",
    ]),
    ("Pipeline \u2014 75 tests", [
        "Scraper Coveo response parsing, research engine (async mocking), scoring deterministic functions",
        "Enricher orchestration, ORM model methods, certification tiers, citation extraction, prompt templates",
    ]),
]
ct = Inches(1.9)
for i, (cat, items) in enumerate(cats):
    y = ct + i * Inches(1.5)
    add_textbox(s, Inches(0.8), y, Inches(11.7), Inches(0.4),
                cat, font_size=19, bold=True)
    for j, item in enumerate(items):
        add_textbox(s, Inches(1.1), y + Inches(0.45) + j * Inches(0.38),
                    Inches(11.0), Inches(0.35),
                    f"\u2022  {item}", font_size=15, color=DARK_GRAY)

add_textbox(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
            "Run:  cd backend && uv run pytest -v",
            font_size=15, color=MED_GRAY, font_name=FONT_MONO)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 15 — TECH STACK
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Tech Stack")

stack = [
    ("Layer", "Technology", "Why"),
    ("Frontend", "Next.js 16 + React 19 + Tailwind 4", "App Router, server components, brand styling"),
    ("Backend", "FastAPI + Python 3.9", "Async-native, auto-generated OpenAPI docs"),
    ("Research LLM", "Perplexity sonar-pro", "Web search with grounded citations"),
    ("Scoring LLM", "OpenAI gpt-4o-mini", "Structured JSON output, lowest cost, fastest"),
    ("Chat Agent", "OpenAI gpt-4o-mini (function calling)", "Agentic loop with 4 domain tools"),
    ("Database", "SQLite (SQLAlchemy ORM)", "Zero-config dev, swap to Postgres with one line"),
    ("Scraping", "Playwright (Chromium)", "Bypasses Akamai, Coveo API interception"),
    ("CLI", "argparse + async", "Independent stages: --force, --limit, --concurrency"),
    ("Testing", "pytest (119 tests)", "In-memory SQLite, async mocking, TestClient"),
]

cols = [Inches(1.5), Inches(4.2), Inches(5.8)]
cxs = [Inches(0.8), Inches(2.3), Inches(6.5)]
ht = Inches(1.8)

for i, (layer, tech, why) in enumerate(stack):
    y = ht + i * Inches(0.52)
    hdr = (i == 0)
    for j, (text, cx, cw) in enumerate(zip([layer, tech, why], cxs, cols)):
        add_textbox(s, cx, y, cw, Inches(0.45),
                    text, font_size=15 if not hdr else 14,
                    bold=hdr, color=BLACK if not hdr else MED_GRAY,
                    font_name=FONT_MONO if hdr else FONT_BODY)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 16 — SCALABILITY & FUTURE
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Scalability: Built to Grow")

items = [
    ("Database \u2192 Postgres", "One connection string. Add Alembic migrations, JSONB with GIN indexes, pgBouncer pooling."),
    ("Pipeline \u2192 Celery + Redis", "Async task queue. Fan out scraping across ZIP codes in parallel workers."),
    ("Search \u2192 Full-Text + Semantic", "Postgres tsvector for name/address. pgvector for semantic search over research summaries."),
    ("Auth \u2192 Multi-Tenant", "JWT auth. Territory assignment \u2014 each rep sees their region. Row-level security."),
    ("Chat \u2192 Streaming", "SSE streaming for real-time token delivery. Multi-turn memory with conversation persistence."),
    ("Observability", "Structured logging, LLM token/latency tracking, pipeline health dashboard with alerting."),
    ("Integrations", "CRM export (Salesforce, HubSpot). Email send from platform. Sync rep notes back."),
]
it = Inches(1.9)
for i, (title, desc) in enumerate(items):
    y = it + i * Inches(0.75)
    add_textbox(s, Inches(0.8), y, Inches(3.2), Inches(0.35),
                title, font_size=17, bold=True)
    add_textbox(s, Inches(4.2), y, Inches(8.5), Inches(0.65),
                desc, font_size=15, color=DARK_GRAY)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 17 — KEY DESIGN DECISIONS (Q&A PREP)
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "Key Design Decisions (Anticipated Q&A)")

decisions = [
    "Why Coveo API interception instead of DOM scraping?\n\u2192  Structured JSON is more reliable, faster, and survives UI redesigns. DOM parsing is fragile.",
    "Why Perplexity + OpenAI instead of one model?\n\u2192  Perplexity excels at grounded web search with citations. OpenAI excels at structured JSON output. Best tool for each job.",
    "Why hybrid scoring instead of fully LLM?\n\u2192  Deterministic where objectively knowable (certs, reviews). LLM where reading comprehension is needed. Reproducible + intelligent.",
    "Why an agentic chat loop instead of pre-canned queries?\n\u2192  Reps ask unpredictable questions. The agent composes tools dynamically \u2014 compare, filter, drill-in \u2014 in ways we can\u2019t anticipate.",
    "Why SQLite for the demo?\n\u2192  Zero-config. SQLAlchemy ORM makes the Postgres swap a one-line change. Schema is already production-ready.",
    "Why sync API handlers when the pipeline is async?\n\u2192  Enricher calls asyncio.run() internally. Async handlers would crash (nested event loops). FastAPI\u2019s thread pool handles it.",
]
add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2),
            decisions, font_size=15, bullet="", spacing=10)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 18 — BY THE NUMBERS
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, OFF_WHITE)
add_subtitle(s)
add_title(s, "By the Numbers")

metrics = [
    ("78", "Contractors scraped\nand enriched"),
    ("119", "Backend tests\npassing"),
    ("26", "Frontend\ncomponents"),
    ("~$0.05", "Total cost to score\nall 78 leads"),
    ("3", "AI enrichment\nstages"),
    ("7", "REST API\nendpoints"),
    ("4", "Agentic chat\ntools"),
    ("5", "Score breakdown\nfactors"),
]
mw, mh, mg = Inches(2.6), Inches(1.6), Inches(0.3)
mx, my = Inches(0.6), Inches(1.9)
for i, (num, label) in enumerate(metrics):
    c, r = i % 4, i // 4
    x = mx + c * (mw + mg)
    y = my + r * (mh + mg)
    add_rect(s, x, y, mw, mh, WHITE, BLACK)
    add_textbox(s, x + Inches(0.15), y + Inches(0.15), mw - Inches(0.3), Inches(0.7),
                num, font_size=36, bold=True, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.15), y + Inches(0.9), mw - Inches(0.3), Inches(0.6),
                label, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 19 — CLOSING
# ─────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BL)
set_slide_bg(s, BLACK)
add_textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
            "Roofing Lead Intelligence Platform",
            font_size=40, bold=True, color=WHITE)
add_bullets(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(3.0), [
    "End-to-End Pipeline \u2014 Scrape, research, score, extract contacts, store, serve, display",
    "Hybrid AI Scoring \u2014 Deterministic + LLM for reproducibility, explainability, and intelligence",
    "Agentic Chat \u2014 Natural language interface with 4 tools and multi-step reasoning over the leads DB",
    "Actionable Intelligence \u2014 Talking points, buying signals, pain points, draft emails, decision-maker contacts",
    "Production-Ready \u2014 119 tests, clean separation, scales to Postgres/Celery/multi-tenant with minimal changes",
], font_size=18, color=LIGHT_GRAY, spacing=14)

add_textbox(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
            "Built by Aryaman Gupta for Instalily AI",
            font_size=16, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════
out = "/Users/aryamangupta/projects-github/instalily-onsite/RoofLeads_Presentation.pptx"
prs.save(out)
print(f"Saved to {out}")
print(f"Total slides: {len(prs.slides)}")
